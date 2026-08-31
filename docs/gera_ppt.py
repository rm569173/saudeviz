"""
Monta o PPTX da Sprint 2 a partir da estrutura definida em estrutura_ppt.md.

Feito por codigo, e nao no PowerPoint, pelo mesmo motivo do diagrama: os
numeros vem do pipeline e mudam. Se a camada Ouro for reprocessada, roda de
novo em vez de caçar slide por slide.

A identidade visual e a do deck da Sprint 1, extraida do arquivo entregue e
documentada em docs/identidade_visual.md: fundo navy escuro, acento teal,
tipografia Cambria.

Saida: EC_Sprint_2_1TSCO_EvidenciasConstrucao_SaudeViz_Lucas_Colen.pptx
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from notas_apresentador import NOTAS

sys.path.insert(0, str(Path(__file__).resolve().parent))

RAIZ = Path(__file__).resolve().parents[1]
IMG = RAIZ / "apresentação"
SAIDA = RAIZ / "EC_Sprint_2_1TSCO_EvidenciasConstrucao_SaudeViz_Lucas_Colen.pptx"

# --- Identidade visual da Sprint 1 ------------------------------------------
FUNDO = RGBColor(0x0A, 0x16, 0x28)
FUNDO_ALT = RGBColor(0x06, 0x10, 0x20)
CARTAO = RGBColor(0x1E, 0x29, 0x3B)
CARTAO_ALTO = RGBColor(0x1E, 0x3A, 0x5F)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_CLARO = RGBColor(0x14, 0xB8, 0xA6)
AMBAR = RGBColor(0xF5, 0x9E, 0x0B)
VERMELHO = RGBColor(0xEF, 0x44, 0x44)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CORPO = RGBColor(0xE2, 0xE8, 0xF0)
FRACO = RGBColor(0x64, 0x74, 0x8B)

SERIF = "Cambria"
MONO = "Consolas"

# Separador de paragrafo dentro das notas do apresentador: linha em branco.
SEPARADOR = chr(10) * 2

# Titulo de cada slide, preenchido durante a montagem e gravado em
# docs/slides.json. O teleprompter le esse manifesto em vez do .pptx: o
# arquivo fica travado enquanto estiver aberto no PowerPoint.
TITULOS: dict[int, str] = {}

# Slide 16:9 no mesmo tamanho do deck da Sprint 1.
LARG, ALT = Inches(10), Inches(5.625)
MARGEM = Inches(0.55)
UTIL = LARG - 2 * MARGEM


def _sem_borda(forma):
    forma.line.fill.background()


def _texto(quadro, linhas, tamanho=11, cor=CORPO, fonte=SERIF, negrito=False,
           espaco=6, alinhamento=PP_ALIGN.LEFT, entrelinha=1.15):
    """Escreve as linhas num text_frame, a primeira reaproveitando o § inicial."""
    quadro.word_wrap = True
    for i, linha in enumerate(linhas):
        p = quadro.paragraphs[0] if i == 0 else quadro.add_paragraph()
        p.alignment = alinhamento
        p.space_after = Pt(espaco)
        p.line_spacing = entrelinha
        # Uma tupla permite dar formato proprio a uma linha isolada.
        conteudo, fmt = (linha, {}) if isinstance(linha, str) else linha
        r = p.add_run()
        r.text = conteudo
        f = r.font
        f.name = fmt.get("fonte", fonte)
        f.size = Pt(fmt.get("tamanho", tamanho))
        f.bold = fmt.get("negrito", negrito)
        f.color.rgb = fmt.get("cor", cor)


def caixa(slide, x, y, w, h, linhas, **kw):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.margin_left = tb.text_frame.margin_right = 0
    tb.text_frame.margin_top = tb.text_frame.margin_bottom = 0
    _texto(tb.text_frame, linhas, **kw)
    return tb


def cartao(slide, x, y, w, h, cor=CARTAO, raio=0.04):
    forma = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    forma.fill.solid()
    forma.fill.fore_color.rgb = cor
    _sem_borda(forma)
    forma.adjustments[0] = raio
    forma.text_frame.text = ""
    return forma


def novo_slide(prs, alt=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fundo = slide.background.fill
    fundo.solid()
    fundo.fore_color.rgb = FUNDO_ALT if alt else FUNDO
    return slide


def titulo(slide, texto, sub=None):
    """Titulo com a barra teal por baixo — assinatura visual do deck."""
    caixa(slide, MARGEM, Inches(0.42), UTIL, Inches(0.5), [texto],
          tamanho=23, cor=BRANCO, negrito=True, espaco=0)
    barra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGEM, Inches(0.95), Inches(0.9), Pt(3.2))
    barra.fill.solid()
    barra.fill.fore_color.rgb = TEAL
    _sem_borda(barra)
    if sub:
        caixa(slide, MARGEM, Inches(1.06), UTIL, Inches(0.3), [sub],
              tamanho=10.5, cor=FRACO, espaco=0)


def anota(slide, numero):
    """
    Grava a fala do apresentador no painel de anotacoes do slide.

    Vai dentro do proprio .pptx: o texto viaja com o arquivo, aparece no modo
    apresentador e nao depende de abrir outro documento na hora de falar.
    """
    nota = NOTAS.get(numero)
    if not nota:
        return
    # As notas sao escritas com quebra de linha por causa da largura do
    # codigo-fonte; aqui viram paragrafos de texto corrido.
    paragrafos = [" ".join(bloco.split())
                  for bloco in nota.split(SEPARADOR) if bloco.strip()]
    quadro = slide.notes_slide.notes_text_frame
    quadro.text = paragrafos[0]
    for paragrafo in paragrafos[1:]:
        quadro.add_paragraph().text = paragrafo


def rodape(slide, numero):
    caixa(slide, MARGEM, ALT - Inches(0.42), Inches(6), Inches(0.25),
          ["SaúdeViz · Challenge FIAP × Oracle 2026 · Sprint 2"],
          tamanho=7.5, cor=FRACO, espaco=0)
    caixa(slide, LARG - MARGEM - Inches(0.6), ALT - Inches(0.42),
          Inches(0.6), Inches(0.25), [str(numero)],
          tamanho=7.5, cor=FRACO, espaco=0, alinhamento=PP_ALIGN.RIGHT)
    anota(slide, numero)


def registra(numero, titulo_slide):
    TITULOS[numero] = titulo_slide


def imagem(slide, nome, x, y, w, h, moldura=True):
    """
    Encaixa a imagem na área preservando a proporção.

    Os prints do painel e do Oracle têm fundo claro; a moldura branca deixa
    explícito que aquilo é captura de outro produto, e não um elemento do
    slide. Sem ela, o print flutua sobre o navy e parece erro de montagem.
    """
    caminho = IMG / nome
    if not caminho.exists():
        cartao(slide, x, y, w, h, CARTAO)
        caixa(slide, x + Inches(0.2), y + h / 2 - Inches(0.15),
              w - Inches(0.4), Inches(0.3), [f"[falta {nome}]"],
              tamanho=10, cor=AMBAR, alinhamento=PP_ALIGN.CENTER)
        return

    with Image.open(caminho) as im:
        pw, ph = im.size
    pad = Emu(int(Inches(0.08))) if moldura else 0
    dw, dh = int(w) - 2 * int(pad), int(h) - 2 * int(pad)
    escala = min(dw / pw, dh / ph)
    fw, fh = int(pw * escala), int(ph * escala)
    fx = int(x) + (int(w) - fw) // 2
    fy = int(y) + (int(h) - fh) // 2

    if moldura:
        fundo = cartao(slide, Emu(fx - int(pad)), Emu(fy - int(pad)),
                       Emu(fw + 2 * int(pad)), Emu(fh + 2 * int(pad)),
                       BRANCO, raio=0.03)
        fundo.shadow.inherit = False
    slide.shapes.add_picture(str(caminho), Emu(fx), Emu(fy),
                             Emu(fw), Emu(fh))


def bloco_dado(slide, x, y, w, h, linhas, tamanho=10):
    """Cartão escuro com texto monoespaçado — números e SQL."""
    cartao(slide, x, y, w, h, CARTAO)
    caixa(slide, x + Inches(0.22), y + Inches(0.16),
          w - Inches(0.44), h - Inches(0.32), linhas,
          tamanho=tamanho, cor=CORPO, fonte=MONO, espaco=2, entrelinha=1.1)


def tabela(slide, x, y, w, cabecalho, linhas, larguras=None, tamanho=9):
    n_l, n_c = len(linhas) + 1, len(cabecalho)
    altura = Inches(0.32) + Inches(0.3) * len(linhas)
    forma = slide.shapes.add_table(n_l, n_c, x, y, w, altura)
    tb = forma.table
    if larguras:
        total = sum(larguras)
        for i, frac in enumerate(larguras):
            tb.columns[i].width = Emu(int(int(w) * frac / total))

    for c, texto in enumerate(cabecalho):
        cel = tb.cell(0, c)
        cel.fill.solid()
        cel.fill.fore_color.rgb = TEAL
        cel.vertical_anchor = MSO_ANCHOR.MIDDLE
        cel.margin_left = cel.margin_right = Inches(0.08)
        _texto(cel.text_frame, [texto], tamanho=tamanho, cor=BRANCO,
               negrito=True, espaco=0)

    for r, linha in enumerate(linhas, start=1):
        for c, texto in enumerate(linha):
            cel = tb.cell(r, c)
            cel.fill.solid()
            cel.fill.fore_color.rgb = CARTAO if r % 2 else CARTAO_ALTO
            cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            cel.margin_left = cel.margin_right = Inches(0.08)
            cor = CORPO
            if texto.startswith("✅"):
                cor = TEAL_CLARO
            elif texto.startswith("⚠️"):
                cor = AMBAR
            elif texto.startswith("🔄"):
                cor = FRACO
            _texto(cel.text_frame, [texto], tamanho=tamanho, cor=cor, espaco=0)
    return forma


# ---------------------------------------------------------------------------
# Composicoes de slide
# ---------------------------------------------------------------------------

CORPO_Y = Inches(1.45)
CORPO_H = Inches(3.6)


def slide_divisor(prs, n, numero, tit, sub):
    registra(n, f"{numero} — {tit}")
    s = novo_slide(prs, alt=True)
    faixa = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.05),
                               Inches(0.16), Inches(1.5))
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = TEAL
    _sem_borda(faixa)
    caixa(s, Inches(0.75), Inches(2.1), UTIL, Inches(0.34), [numero],
          tamanho=12, cor=TEAL_CLARO, negrito=True, espaco=4)
    caixa(s, Inches(0.75), Inches(2.5), UTIL, Inches(0.6), [tit],
          tamanho=30, cor=BRANCO, negrito=True, espaco=6)
    caixa(s, Inches(0.75), Inches(3.15), Inches(7.5), Inches(0.5), [sub],
          tamanho=11, cor=FRACO)
    rodape(s, n)
    return s


def slide_img(prs, n, tit, sub, imgs, nota=None):
    registra(n, tit)
    s = novo_slide(prs)
    titulo(s, tit, sub)
    h = CORPO_H - (Inches(0.55) if nota else 0)
    if len(imgs) == 1:
        imagem(s, imgs[0], MARGEM, CORPO_Y, UTIL, h)
    else:
        gap = Inches(0.2)
        w = Emu((int(UTIL) - int(gap) * (len(imgs) - 1)) // len(imgs))
        for i, nome in enumerate(imgs):
            imagem(s, nome, Emu(int(MARGEM) + i * (int(w) + int(gap))),
                   CORPO_Y, w, h)
    if nota:
        caixa(s, MARGEM, CORPO_Y + h + Inches(0.12), UTIL, Inches(0.45),
              [nota], tamanho=9.5, cor=FRACO)
    rodape(s, n)
    return s


def _altura_bloco(linhas, tamanho):
    """Altura que o cartão precisa para caber o texto, sem sobra."""
    # 1,58 e o passo real que Consolas ocupa com entrelinha 1,1 mais o
    # espaco entre paragrafos; medido no render, nao no nominal.
    return Emu(int(Inches(0.38) + Pt(tamanho * 1.58) * len(linhas)))


def slide_img_bloco(prs, n, tit, sub, img, bloco, nota=None, larg_img=None,
                    tam_bloco=9):
    """
    Print à esquerda, dados à direita.

    A largura da coluna da imagem se ajusta à proporção do print: as capturas
    do painel são muito largas (2:1 ou mais) e, numa coluna estreita, o texto
    dentro delas fica ilegível. O cartão de dados tem a altura do próprio
    conteúdo e fica centrado — esticá-lo até o rodapé deixava metade vazia.
    """
    registra(n, tit)
    s = novo_slide(prs)
    titulo(s, tit, sub)
    gap = Inches(0.22)
    h = CORPO_H - (Inches(0.5) if nota else 0)

    if larg_img is None:
        caminho = IMG / img
        if caminho.exists():
            with Image.open(caminho) as im:
                proporcao = im.size[0] / im.size[1]
            larg_img = 0.72 if proporcao >= 1.9 else 0.6
        else:
            larg_img = 0.6

    wi = Emu(int((int(UTIL) - int(gap)) * larg_img))
    wb = Emu(int(UTIL) - int(gap) - int(wi))
    xb = Emu(int(MARGEM) + int(wi) + int(gap))

    imagem(s, img, MARGEM, CORPO_Y, wi, h)
    hb = min(int(_altura_bloco(bloco, tam_bloco)), int(h))
    yb = Emu(int(CORPO_Y) + (int(h) - hb) // 2)
    bloco_dado(s, xb, yb, wb, Emu(hb), bloco, tamanho=tam_bloco)

    if nota:
        caixa(s, MARGEM, CORPO_Y + h + Inches(0.1), UTIL, Inches(0.4),
              [nota], tamanho=9.5, cor=FRACO)
    rodape(s, n)
    return s


def slide_bloco(prs, n, tit, sub, bloco, nota=None, tam=10.5, larg=0.8):
    registra(n, tit)
    s = novo_slide(prs)
    titulo(s, tit, sub)
    w = Emu(int(int(UTIL) * larg))
    disponivel = CORPO_H - (Inches(0.6) if nota else 0)
    h = Emu(min(int(_altura_bloco(bloco, tam)), int(disponivel)))
    bloco_dado(s, MARGEM, CORPO_Y, w, h, bloco, tamanho=tam)
    if nota:
        caixa(s, MARGEM, CORPO_Y + h + Inches(0.14), UTIL, Inches(0.5),
              [nota], tamanho=9.5, cor=FRACO)
    rodape(s, n)
    return s


def slide_tabela(prs, n, tit, sub, cab, linhas, larguras=None, tam=9,
                 nota=None):
    registra(n, tit)
    s = novo_slide(prs)
    titulo(s, tit, sub)
    tabela(s, MARGEM, CORPO_Y, UTIL, cab, linhas, larguras, tam)
    if nota:
        caixa(s, MARGEM, ALT - Inches(0.85), UTIL, Inches(0.4), [nota],
              tamanho=9, cor=FRACO)
    rodape(s, n)
    return s


# ---------------------------------------------------------------------------
# Montagem
# ---------------------------------------------------------------------------

def capa(prs, n):
    s = novo_slide(prs, alt=True)
    barra = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), ALT)
    barra.fill.solid()
    barra.fill.fore_color.rgb = TEAL
    _sem_borda(barra)
    caixa(s, Inches(0.95), Inches(1.55), Inches(8), Inches(1.0), ["SaúdeViz"],
          tamanho=52, cor=BRANCO, negrito=True, espaco=4)
    caixa(s, Inches(0.95), Inches(2.62), Inches(8), Inches(0.5),
          ["Painel inteligente de acesso hospitalar e perfil de atendimento"],
          tamanho=15, cor=TEAL_CLARO, espaco=4)
    risco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95),
                               Inches(3.25), Inches(1.4), Pt(2.5))
    risco.fill.solid()
    risco.fill.fore_color.rgb = AMBAR
    _sem_borda(risco)
    caixa(s, Inches(0.95), Inches(3.55), Inches(8), Inches(1.1), [
        "Challenge FIAP × Oracle 2026 · Sprint 2",
        "Turma 1TSCOA",
        "Lucas Ventura Araujo Ribas Colen — RM 569173",
    ], tamanho=11.5, cor=CORPO, espaco=4)
    # A capa nao leva rodape, mas consome numero como qualquer slide:
    # sem isso as notas ficam deslocadas em um.
    registra(n, "SaúdeViz — capa")
    anota(s, n)
    return s


def identificacao(prs, n):
    registra(n, "Identificação do grupo")
    s = novo_slide(prs)
    titulo(s, "Identificação do grupo")
    campos = [
        ("INTEGRANTE DO GRUPO", "Lucas Ventura Araujo Ribas Colen"),
        ("RM", "569173"),
        ("TURMA", "1TSCOA"),
        ("INSTITUIÇÃO", "FIAP"),
        ("ANO", "2026"),
    ]
    y = CORPO_Y + Inches(0.1)
    for rotulo, valor in campos:
        caixa(s, MARGEM, y, Inches(2.6), Inches(0.28), [rotulo],
              tamanho=9, cor=TEAL_CLARO, negrito=True, espaco=0)
        caixa(s, Inches(3.3), y - Inches(0.03), Inches(5.8), Inches(0.32),
              [valor], tamanho=13, cor=BRANCO, espaco=0)
        y += Inches(0.62)
    rodape(s, n)
    return s


def numeros(prs, n):
    registra(n, "Os números da entrega")
    s = novo_slide(prs)
    titulo(s, "Os números da entrega", "Camada Ouro carregada no Oracle 19c "
                                       "da FIAP")
    itens = [
        ("5.546.817", "internações do SIH/SUS\nocorridas em 2024"),
        ("R$ 10,03 bi", "pagos pelo SUS\nno período"),
        ("4 estados", "ES · MG · RJ · SP\n89 milhões de habitantes"),
        ("3.131", "estabelecimentos\ncom leito"),
        ("11", "tabelas na camada\nGold do Oracle"),
    ]
    gap = Inches(0.14)
    w = Emu((int(UTIL) - int(gap) * 4) // 5)
    for i, (valor, rotulo) in enumerate(itens):
        x = Emu(int(MARGEM) + i * (int(w) + int(gap)))
        cartao(s, x, CORPO_Y, w, Inches(1.55))
        caixa(s, x + Inches(0.14), CORPO_Y + Inches(0.24),
              w - Inches(0.28), Inches(0.45), [valor],
              tamanho=17, cor=TEAL_CLARO, negrito=True, espaco=2,
              alinhamento=PP_ALIGN.CENTER)
        caixa(s, x + Inches(0.12), CORPO_Y + Inches(0.78),
              w - Inches(0.24), Inches(0.65), rotulo.split("\n"),
              tamanho=8.5, cor=FRACO, espaco=1,
              alinhamento=PP_ALIGN.CENTER, entrelinha=1.1)
    imagem(s, "eda_resumo_executivo.png", MARGEM, CORPO_Y + Inches(1.72),
           UTIL, Inches(1.75))
    rodape(s, n)
    return s


def encerramento(prs, n):
    registra(n, "Encerramento")
    s = novo_slide(prs, alt=True)
    barra = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), ALT)
    barra.fill.solid()
    barra.fill.fore_color.rgb = TEAL
    _sem_borda(barra)
    caixa(s, Inches(0.95), Inches(1.75), Inches(8), Inches(1.1), [
        "Dados que salvam vidas.",
        "Decisões que transformam o sistema de saúde.",
    ], tamanho=26, cor=BRANCO, negrito=True, espaco=8)
    risco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95),
                               Inches(3.25), Inches(1.4), Pt(2.5))
    risco.fill.solid()
    risco.fill.fore_color.rgb = AMBAR
    _sem_borda(risco)
    caixa(s, Inches(0.95), Inches(3.6), Inches(8), Inches(0.8), [
        "saudeviz.streamlit.app",
        "github.com/rm569173/saudeviz",
    ], tamanho=13, cor=TEAL_CLARO, espaco=5)
    rodape(s, n)
    return s


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = LARG, ALT
    n = 0

    def prox():
        nonlocal n
        n += 1
        return n

    capa(prs, prox())
    identificacao(prs, prox())
    numeros(prs, prox())

    # --- 1a entrega ---------------------------------------------------------
    slide_divisor(prs, prox(), "1ª ENTREGA", "Sprint 1 atualizada",
                  "O que foi prometido, o que foi entregue e o que mudou "
                  "de rota — com o motivo.")
    slide_tabela(
        prs, prox(), "O que mudou desde a Sprint 1", None,
        ["Prometido na Sprint 1", "Situação"],
        [["3 fontes: SQL, JSON e CSV", "✅ Entregue — e uma quarta, de clima"],
         ["Medallion Bronze → Prata → Ouro",
          "✅ Entregue, em PySpark no Databricks"],
         ["Oracle como camada Gold", "✅ Entregue, 11 tabelas T_SAUDE_*"],
         ["Análise preditiva", "✅ Entregue, com validação temporal"],
         ["Oracle Select AI",
          "⚠️ Indisponível no 19c — equivalente implementado"],
         ["External Table do CSV",
          "⚠️ Sem privilégio na conta acadêmica — DDL entregue"],
         ["Power BI", "🔄 Substituído por Streamlit"],
         ["Cobertura nacional", "🔄 Recorte no Sudeste para o MVP"]],
        larguras=[0.42, 0.58], tam=9.5)
    slide_img_bloco(
        prs, prox(), "Quadro de gestão",
        "Kanban público, com o que ficou por fazer no status real",
        "trello.png",
        ["QUADRO", "trello.com/b/3XTAInGQ", "",
         "44 cartões · 5 listas", "",
         "Sprint 1              8", "Concluído           23",
         "Em andamento         2", "A fazer              5",
         "Não implementado     5", "",
         "A quinta lista é a que", "mostra decisão de escopo",
         "em vez de esquecimento."],
        larg_img=0.62, tam_bloco=8.5)

    # --- 2a entrega ---------------------------------------------------------
    slide_divisor(prs, prox(), "2ª ENTREGA", "MVP implementado",
                  "Painel público conectado ao Oracle da FIAP em tempo real.")
    slide_img(prs, prox(), "O painel", "saudeviz.streamlit.app",
              ["painel_01_visao_geral.png"],
              "Cinco páginas. O selo na barra lateral mostra se o dado veio "
              "do Oracle ou do retrato local — resultado de contingência "
              "nunca se passa por dado ao vivo.")
    slide_img_bloco(
        prs, prox(), "Onde a capacidade foi ultrapassada",
        "Classificação por município e mês",
        "painel_02_capacidade_status.png",
        ["Folga           8.466", "Adequada        1.043",
         "Atenção           223", "Crítica            16", "",
         "município-mês, em 823", "municípios com internação",
         "registrada em 2024"])
    slide_img(prs, prox(), "Municípios sob pressão", None,
              ["painel_03_capacidade_criticos.png"],
              "Ocupação acima de 1,0 é alerta para investigar, não prova de "
              "colapso: pode ser sobrecarga real, leito desatualizado no CNES "
              "ou município-polo atendendo toda uma região.")
    slide_img_bloco(
        prs, prox(), "Quais perfis pressionam mais o sistema",
        "Pressão relativa: participação nos leitos-dia ÷ participação nas "
        "internações",
        "painel_04_perfis_pressao.png",
        ["TRANSTORNOS MENTAIS", "",
         "1,9% das internações", "4,1% dos leitos-dia", "",
         "pressão relativa  2,15", "permanência      10,4 d",
         "custo médio      R$ 521", "",
         "Ocupa leito e não consome", "orçamento. Invisível num",
         "painel que só conta", "atendimentos."],
        larg_img=0.62, tam_bloco=8.5)

    # --- 3a entrega ---------------------------------------------------------
    slide_divisor(prs, prox(), "3ª ENTREGA", "Arquitetura final",
                  "Quatro fontes públicas, Medallion no Databricks, "
                  "camada de consumo no Oracle.")
    s = novo_slide(prs)
    imagem(s, "arquitetura_solucao.png", Inches(0.15), Inches(0.15),
           LARG - Inches(0.3), ALT - Inches(0.3), moldura=False)
    n_arq = prox()
    registra(n_arq, "Desenho da arquitetura")
    rodape(s, n_arq)
    slide_tabela(
        prs, prox(), "O que roda onde, e por quê", None,
        ["Etapa", "Onde", "Motivo"],
        [["Download e decode .dbc", "Estação local",
          "Formato proprietário do DATASUS, exige extensão C e FTP"],
         ["Bronze, Prata, Ouro", "Databricks",
          "PySpark sobre tabelas Delta"],
         ["Camada Gold", "Oracle 19c FIAP",
          "Camada de serviço consultada pelo painel"],
         ["Painel e NL→SQL", "Streamlit Cloud",
          "Link público, código versionado"]],
        larguras=[0.26, 0.2, 0.54], tam=9.5)
    slide_img(prs, prox(), "Medallion materializada",
              "Catálogo do Unity Catalog, não desenho",
              ["databricks_catalog_medallion.png"],
              "As três camadas existem como schemas reais no Databricks.")
    slide_img(prs, prox(), "A camada Gold no Oracle", None,
              ["oracle_contagem_tabelas.png", "oracle_arvore_tabelas.png"])
    slide_img_bloco(
        prs, prox(), "Integridade do modelo",
        "Três tabelas construídas por caminhos diferentes",
        "oracle_integridade_totais.png",
        ["Fato de internações", "         5.546.817", "",
         "Indicador de capacidade", "         5.546.817", "",
         "Ranking de hospitais", "         5.546.817", "",
         "Mesmo total.", "Nenhuma internação",
         "perdida ou duplicada."], tam_bloco=9)
    slide_tabela(
        prs, prox(), "O que não foi implementado, e por quê",
        "Decisão registrada, não omissão",
        ["Item", "Motivo", "O que foi entregue"],
        [["Oracle Select AI",
          "Exige Autonomous. Verificado em all_objects: o 19c da FIAP não "
          "tem DBMS_CLOUD_AI",
          "Tradutor NL→SQL sobre os mesmos metadados + script pronto"],
         ["External Table", "Conta acadêmica sem CREATE ANY DIRECTORY",
          "CSV como tabela comum; DDL documentado"],
         ["Cobertura nacional", "Recorte no Sudeste para viabilizar o MVP",
          "Pipeline parametrizado por UF"],
         ["Clusterização K-Means", "Cortada por prazo",
          "Código de referência no repositório"]],
        larguras=[0.2, 0.42, 0.38], tam=8.5)

    # --- 4a entrega ---------------------------------------------------------
    slide_divisor(prs, prox(), "4ª ENTREGA", "Modelos e técnicas",
                  "Ingestão das fontes, a descoberta que reorientou o "
                  "projeto, estatística e previsão.")
    slide_bloco(
        prs, prox(), "As fontes e seus formatos",
        "Três formatos exigidos pelo desafio, mais uma quarta fonte",
        [("FONTE 1 — SIH/SUS       relacional (.dbc)", {"cor": TEAL_CLARO}),
         "  ftp.retrbinary(...)              7.015.106 internações", "",
         ("FONTE 2 — CNES          JSON via API REST", {"cor": TEAL_CLARO}),
         "  _get_json(...)                   4.481 estabelecimentos", "",
         ("FONTE 3 — IBGE          CSV", {"cor": TEAL_CLARO}),
         "  df.to_csv(...)                   5.571 municípios", "",
         ("FONTE 4 — Open-Meteo    JSON via API REST", {"cor": AMBAR}),
         "  archive-api.open-meteo.com       366 dias × 4 capitais"],
        nota="A quarta fonte não era exigida. Entrou para testar hipóteses "
             "que os dados de saúde sozinhos não respondem.",
        tam=10, larg=0.92)
    slide_img(prs, prox(), "Tratamento do formato proprietário",
              "O .dbc do DATASUS é um DBF comprimido com PKWare DCL",
              ["extrai_sih_py_69.jpeg"],
              "A conversão roda em diretório temporário com caminho ASCII: a "
              "extensão C não aceita acento no caminho, e o projeto mora em "
              "\"Educação/Ciência de Dados\".")
    slide_img_bloco(
        prs, prox(), "A descoberta que reorientou o projeto",
        "Defasagem entre internar e ser faturado",
        "pipeline_prata_defasagem.png",
        ["0 meses  4.214.705  60,8%", "1 mês    1.789.226  25,8%",
         "2 meses    587.295   8,5%", "3 meses    297.454   4,3%", "",
         ("acumulado em M+3: 99,4%", {"cor": TEAL_CLARO}), "",
         "A competência do SIH é o", "mês de PAGAMENTO, não o",
         "da internação.", "",
         "Corrigimos a dimensão", "temporal e ampliamos a",
         "ingestão até março/2025."], tam_bloco=8.5)
    slide_img(prs, prox(), "Análise exploratória",
              "16 consultas SQL cobrindo as quatro frentes do desafio",
              ["eda_q03_sazonalidade_grafico.png",
               "eda_q07_pressao_grafico.png"])
    slide_img_bloco(
        prs, prox(), "Técnicas estatísticas",
        "Outliers por IQR de Tukey e correlação de Pearson",
        "eda_q14_outliers_iqr.png",
        ["TUKEY (Q3 + 1,5×IQR)", "limite de outlier: 1,172", "",
         "O critério estatístico,", "calculado sem referência",
         "ao nosso limiar, confirma", "que ocupação acima de 1,0",
         "é anômala nesta", "distribuição.", "",
         ("PEARSON", {"cor": TEAL_CLARO}),
         "ocupação × transferência", "        −0,305", "",
         "Municípios com baixa", "ocupação transferem mais:",
         "estabilizam e encaminham."], larg_img=0.55, tam_bloco=8)
    slide_img_bloco(
        prs, prox(), "Modelo de previsão de demanda",
        "Validação por janela expansível — o modelo mais simples venceu",
        "modelo_comparativo_horizonte.png",
        ["horizonte  semanal  regressão", "",
         "      7 d    6,32%     6,27%", "     30 d    5,77%     7,67%",
         ("     90 d    5,46%    27,12%", {"cor": AMBAR}), "",
         "Testamos regressão com", "tendência e sazonalidade.",
         "Perdeu em todos os", "horizontes acima de 7 dias.", "",
         "A série não tem tendência", "explorável."], tam_bloco=8.5)
    slide_img_bloco(
        prs, prox(), "O que move a demanda hospitalar",
        "Variação percentual frente à segunda-feira",
        "modelo_dia_semana.png",
        [("Sábado    −38%", {"cor": AMBAR}),
         ("Domingo   −41%", {"cor": AMBAR}),
         ("Feriado   −26%", {"cor": AMBAR}), "",
         "Consistente nos quatro", "estados.", "",
         "A queda do fim de semana", "não é falta de doente:",
         "é a rede eletiva parada."], tam_bloco=9)

    # --- 5a entrega ---------------------------------------------------------
    slide_divisor(prs, prox(), "5ª ENTREGA", "Evidências visuais",
                  "O painel e a camada Gold respondendo às perguntas "
                  "do desafio.")
    slide_img(prs, prox(), "Previsão de demanda no painel",
              "Projeção diária de 90 dias, com intervalo de 95%",
              ["painel_05_previsao_serie.png"])
    slide_img_bloco(
        prs, prox(), "Quantos leitos, e quando abrir",
        "Consumo observado traduzido em dimensionamento",
        "painel_07_dimensionamento.png",
        [("BELO HORIZONTE", {"cor": TEAL_CLARO}),
         "6.312 leitos cadastrados", "",
         "Demanda comum      5.107", "  abertos o ano inteiro", "",
         "Pico (abril)       5.904", "  +797 leitos sazonais", "",
         ("Folga no pico        408", {"cor": AMBAR}),
         ("  6% da capacidade", {"cor": AMBAR}), "",
         "Abril é o pico nas", "quatro capitais."],
        nota="Seis por cento de folga no pico é a margem que uma epidemia "
             "consome em dias.", larg_img=0.6, tam_bloco=8.5)
    slide_img_bloco(
        prs, prox(), "Consultas na camada Gold",
        "Executadas no Oracle 19c da FIAP",
        "oracle_top10_hospitais.png",
        ["Santa Casa BH", "  54.406 internações",
         "  0,21% transferência", "",
         "HC-FMUSP SP", "  53.161 · 2,66%", "",
         "Hosp. Base SJRP", "  50.043 · 1,09%"], tam_bloco=9)
    slide_img(prs, prox(), "Onde as internações crescem",
              "Segundo semestre contra o primeiro, por município",
              ["eda_q05_crescimento.png"])
    slide_img_bloco(
        prs, prox(), "Quem exporta paciente",
        "A taxa de transferência mede resolutividade, não volume",
        "eda_q12_transferencias.png",
        [("EMBU-GUAÇU / SP", {"cor": TEAL_CLARO}), "",
         "1.460 internações", "1.082 transferências",
         "   15 leitos", "",
         ("74% dos pacientes seguem", {"cor": AMBAR}),
         ("para outro município.", {"cor": AMBAR}), "",
         "Não é \"há muitas", "internações aqui\": é",
         "\"pacientes saem daqui", "porque não há como",
         "tratá-los aqui\"."], tam_bloco=8.5)

    # --- Analise complementar: clima ---------------------------------------
    slide_divisor(prs, prox(), "ANÁLISE COMPLEMENTAR", "Clima e internação",
                  "Não é entrega obrigatória. Três hipóteses testadas — "
                  "uma refutada, uma confirmada, uma com ressalva.")
    slide_img_bloco(
        prs, prox(), "Uma quarta fonte de dados",
        "Open-Meteo · clima diário das 4 capitais · 2024",
        "clima_01_fonte.png",
        ["1.464 dias-capital", "  526 com chuva", "",
         "1.421.225 internações", "nas capitais", "",
         ("TRÊS HIPÓTESES", {"cor": TEAL_CLARO}), "",
         "chuva aumenta acidentes?", "",
         "frio aumenta doença", "respiratória?", "",
         "a estação importa além", "do ciclo semanal?"], tam_bloco=8.5)
    slide_img_bloco(
        prs, prox(), "Chuva não explica acidentes",
        "Hipótese refutada — e o teste de gradiente é o que prova",
        "clima_03_gradiente_chuva.png",
        ["Com chuva   104,4 /dia", "Sem chuva   105,1",
         ("variação     −0,7%", {"cor": VERMELHO}), "",
         ("SEM GRADIENTE", {"cor": TEAL_CLARO}), "",
         "Sem chuva      693 d  105,4", "Chuva fraca    254    103,0",
         "Chuva moderada  85    110,4", "Chuva forte     16     83,1"],
        nota="Se a chuva causasse acidentes, o efeito cresceria com a "
             "intensidade. Sobe e desce sem padrão — e os 83,1 da chuva "
             "forte vêm de 16 dias apenas. O planejamento hospitalar não "
             "deve reservar capacidade para dias de chuva.",
        larg_img=0.55, tam_bloco=8)
    slide_img_bloco(
        prs, prox(), "Frio se associa a internação respiratória",
        "Hipótese confirmada — gradiente monotônico nas quatro capitais",
        "clima_04_frio_respiratoria.png",
        ["Muito frio  73,7 /dia", "Quente      67,5",
         ("variação    +9,2%", {"cor": TEAL_CLARO}), "",
         "Capital      MFrio  Quente", "BH            61,9    50,0",
         "RJ            50,5    43,6", "SP           175,1   163,9"],
        nota="As faixas são quartis calculados dentro de cada capital: "
             "18 °C é frio em Vitória e ameno em São Paulo. Um limiar "
             "absoluto compararia climas diferentes como se fossem o mesmo.",
        larg_img=0.55, tam_bloco=8.5)
    slide_img_bloco(
        prs, prox(), "Só a temperatura mínima importa",
        "Correlação de Pearson com internação respiratória",
        "clima_05_correlacao.png",
        ["           tmin   tmax  chuva", "",
         "BH       −0,306  0,020 −0,244", "RJ       −0,185 −0,001 −0,129",
         "Vitória  −0,148 −0,055 −0,181", "SP       −0,107  0,062 −0,139"],
        nota="A temperatura máxima não tem correlação — ronda zero nas "
             "quatro cidades. É o frio da madrugada que se relaciona com "
             "internação respiratória, não o calor do dia.",
        larg_img=0.55, tam_bloco=8.5)
    slide_bloco(
        prs, prox(), "A ressalva que o próprio dado revela",
        "Por que isto é associação, e não causa",
        [("A chuva aparece com correlação NEGATIVA:", {"cor": AMBAR}),
         ("dias chuvosos têm menos internação respiratória.",
          {"cor": AMBAR}), "",
         "Contraintuitivo — até lembrar que no Sudeste chove no verão",
         "e a seca é no inverno. A chuva está medindo a ESTAÇÃO,",
         "não a si mesma.", "",
         "O mesmo confundimento pode valer para o frio: dias frios",
         "concentram-se no inverno, que também tem mais circulação",
         "viral. Os 9,2% podem ser efeito da temperatura, da",
         "sazonalidade viral, ou dos dois.", "",
         ("É associação, não causa.", {"cor": TEAL_CLARO})],
        nota="Separar os dois efeitos exigiria dados de circulação viral "
             "por semana epidemiológica — fora do escopo desta sprint, e "
             "registrado como próximo passo.", tam=9.5, larg=0.95)

    # --- Select AI ----------------------------------------------------------
    slide_divisor(prs, prox(), "SELECT AI", "Perguntas em português",
                  "O mecanismo equivalente, sobre os mesmos metadados que o "
                  "Select AI usaria.")
    slide_img(prs, prox(), "Perguntas em linguagem natural",
              'Página "Pergunte em português" do painel',
              ["painel_08_nlsql_exemplos.png"])
    slide_img_bloco(
        prs, prox(), "O SQL gerado e executado",
        "Duas perguntas, duas tabelas diferentes",
        "painel_10_nlsql_sql.png",
        ['"quais hospitais mais', ' transferem pacientes?"',
         ("  → T_SAUDE_RANK_HOSPITAIS", {"cor": TEAL_CLARO}), "",
         '"onde a capacidade está', ' sendo ultrapassada?"',
         ("  → T_SAUDE_IND_", {"cor": TEAL_CLARO}),
         ("    CAPACIDADE_MUNICIPAL", {"cor": TEAL_CLARO}), "",
         "A tabela é escolhida pela", "intenção reconhecida, não",
         "por uma consulta fixa."],
        nota="O SQL continua sendo gerado e executado no Oracle. O tradutor "
             "remove a barreira da sintaxe para quem decide.",
        larg_img=0.6, tam_bloco=8.5)
    slide_img_bloco(
        prs, prox(), "Por que não é o Select AI nativo",
        "Verificado no banco, não presumido",
        "oracle_comentarios_metadados.png",
        ["SELECT COUNT(*)", "  FROM all_objects",
         " WHERE object_name =", "       'DBMS_CLOUD_AI';", "",
         ("            →  0", {"cor": VERMELHO}), "",
         "O Select AI existe apenas", "no Autonomous Database.", "",
         "Implementamos o mecanismo", "equivalente sobre OS MESMOS",
         "metadados: os COMMENT ON", "das tabelas."],
        nota="Migrar para o Autonomous não exige remodelar nada — só trocar "
             "o motor de tradução. O script já está entregue.",
        larg_img=0.58, tam_bloco=8.5)

    # --- Fechamento: 6a, 7a e 8a entrega ------------------------------------
    slide_divisor(prs, prox(), "6ª · 7ª · 8ª ENTREGA", "Fechamento",
                  "Repositório, vídeo pitch e conclusões.")
    slide_img_bloco(
        prs, prox(), "Código-fonte", "github.com/rm569173/saudeviz",
        "github_repositorio.png",
        ["6 notebooks Databricks", "",
         "pipeline de ingestão das", "quatro fontes", "",
         "DDL do Oracle com", "COMMENT ON", "",
         "tradutor NL→SQL", "", "painel Streamlit", "",
         "camada Gold versionada", "para reprodução sem", "acesso ao banco"],
        larg_img=0.62, tam_bloco=8.5)
    slide_bloco(
        prs, prox(), "Vídeo pitch", "Até 5 minutos",
        [("youtube.com/watch?v=________", {"cor": TEAL_CLARO}), "",
         "Preencher o link após o upload.", "",
         "Publicar como NÃO LISTADO ou público — vídeo privado",
         "não abre para o avaliador."], tam=11, larg=0.75)
    slide_bloco(
        prs, prox(), "O que aprendemos com os dados", None,
        [("1. A competência do SIH não é a data da internação",
          {"cor": TEAL_CLARO}),
         "   42% dos registros de um mês são de meses anteriores.",
         "   Sem corrigir isso, toda série temporal do projeto",
         "   estaria errada.", "",
         ("2. A coluna COBRANCA revela transferência de paciente",
          {"cor": TEAL_CLARO}),
         '   Muda "há muitas internações aqui" para "pacientes saem',
         '   daqui porque não há como tratá-los aqui".', "",
         ("3. A demanda hospitalar não tem tendência explorável",
          {"cor": TEAL_CLARO}),
         "   O sinal previsível está no ciclo semanal e nos feriados.",
         "   O modelo simples venceu o sofisticado por 5x em 90 dias."],
        tam=9.5, larg=0.95)
    slide_bloco(
        prs, prox(), "Limitações declaradas",
        "O que este painel não pode afirmar",
        ["·  Recorte no Sudeste — o pipeline é parametrizado por UF",
         "",
         "·  Dezembro/2024 com cobertura de 99,4%, não 100%",
         "",
         "·  Ocupação acima de 1,0 é alerta, não diagnóstico de colapso",
         "",
         "·  Pacientes internados antes de 2024 não entram na contagem",
         "   (medido: menos de 1% dos leitos-dia)",
         "",
         "·  O efeito do frio é associação, não causa comprovada"],
        tam=10, larg=0.9)
    slide_bloco(
        prs, prox(), "Próximos passos", None,
        ["·  Separar o efeito do frio do efeito da sazonalidade viral,",
         "   cruzando com circulação respiratória por semana",
         "   epidemiológica",
         "",
         "·  Estender ao Brasil inteiro — o pipeline já é parametrizado",
         "   por UF, é trocar uma lista em config.py",
         "",
         "·  Migrar para o Autonomous Database e ativar o Select AI",
         "   nativo — o script de configuração já está entregue",
         "",
         "·  Retomar a clusterização de municípios por perfil"],
        tam=10, larg=0.9)
    encerramento(prs, prox())

    # O manifesto sai antes do deck: se o .pptx estiver aberto no PowerPoint
    # a gravacao dele falha, e nao ha razao para perder o manifesto junto.
    manifesto = Path(__file__).resolve().parent / "slides.json"
    manifesto.write_text(
        json.dumps({str(k): v for k, v in sorted(TITULOS.items())},
                   ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"manifesto: {manifesto.name} ({len(TITULOS)} títulos)")

    try:
        prs.save(SAIDA)
    except PermissionError:
        raise SystemExit(
            f"{SAIDA.name} esta aberto em outro programa e nao pode ser "
            f"gravado. Feche o PowerPoint e rode de novo.")

    print(f"gravado: {SAIDA.name}")
    print(f"{n} slides · {SAIDA.stat().st_size / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    main()
